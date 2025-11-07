"""
Slide template layouts for Basin Climbing presentations.

Provides reusable slide layout methods with consistent branding.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import List, Dict, Optional, Tuple
from PIL import Image as PILImage
import pandas as pd


# Basin Climbing Brand Colors
COLORS = {
    'primary_blue': RGBColor(31, 119, 180),    # #1f77b4
    'teal': RGBColor(44, 127, 184),            # #2c7fb8
    'orange': RGBColor(255, 127, 80),          # #ff7f50
    'green': RGBColor(44, 160, 44),            # #2ca02c
    'yellow': RGBColor(255, 193, 7),           # #ffc107
    'gray': RGBColor(127, 127, 127),           # #7f7f7f
    'dark_gray': RGBColor(44, 62, 80),         # #2c3e50
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(0, 0, 0),
}


class SlideTemplate:
    """
    Provides methods to create consistently-styled slides for Basin Climbing presentations.
    """

    def __init__(self, presentation: Presentation):
        """
        Initialize with a python-pptx Presentation object.

        Args:
            presentation: pptx.Presentation instance
        """
        self.prs = presentation
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

    def _add_blank_slide(self):
        """Add a blank slide and return it."""
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        return self.prs.slides.add_slide(blank_layout)

    def _add_title(self, slide, title: str, top: Inches = Inches(0.5),
                   font_size: int = 32, color=COLORS['primary_blue']):
        """Add a title text box to a slide."""
        left = Inches(0.5)
        width = self.slide_width - Inches(1)
        height = Inches(0.8)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title

        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = True
        paragraph.font.color.rgb = color
        paragraph.font.name = 'Calibri'

        return title_box

    def _add_footer(self, slide, text: str = "Basin Climbing & Fitness"):
        """Add footer text to a slide."""
        left = Inches(0.5)
        top = self.slide_height - Inches(0.5)
        width = self.slide_width - Inches(1)
        height = Inches(0.3)

        footer_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = footer_box.text_frame
        text_frame.text = text

        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = COLORS['gray']
        paragraph.font.name = 'Calibri'

    def _add_text_box(self, slide, text: str, left: Inches, top: Inches,
                      width: Inches, height: Inches, font_size: int = 18,
                      bold: bool = False, color=COLORS['black'], align=PP_ALIGN.LEFT):
        """Add a text box with specified styling."""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True

        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = 'Calibri'
        paragraph.alignment = align

        return text_box

    def title_slide(self, title: str, subtitle: str, date: str = "") -> None:
        """
        Create a title slide.

        Args:
            title: Main title text
            subtitle: Subtitle text
            date: Date or date range
        """
        slide = self._add_blank_slide()

        # Main title - centered
        title_left = Inches(1)
        title_top = self.slide_height / 2 - Inches(1)
        title_width = self.slide_width - Inches(2)
        title_height = Inches(1.5)

        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        text_frame = title_box.text_frame
        text_frame.text = title

        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLORS['primary_blue']
        paragraph.font.name = 'Calibri'
        paragraph.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_top = title_top + Inches(1.5)
        self._add_text_box(
            slide, subtitle,
            title_left, subtitle_top, title_width, Inches(0.8),
            font_size=24, color=COLORS['dark_gray'], align=PP_ALIGN.CENTER
        )

        # Date
        if date:
            date_top = subtitle_top + Inches(0.8)
            self._add_text_box(
                slide, date,
                title_left, date_top, title_width, Inches(0.5),
                font_size=18, color=COLORS['gray'], align=PP_ALIGN.CENTER
            )

    def section_header(self, text: str) -> None:
        """
        Create a section header slide.

        Args:
            text: Section title text
        """
        slide = self._add_blank_slide()

        # Centered section title
        title_left = Inches(1)
        title_top = self.slide_height / 2 - Inches(0.5)
        title_width = self.slide_width - Inches(2)
        title_height = Inches(1)

        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        text_frame = title_box.text_frame
        text_frame.text = text

        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLORS['primary_blue']
        paragraph.font.name = 'Calibri'
        paragraph.alignment = PP_ALIGN.CENTER

        self._add_footer(slide)

    def metric_cards_slide(self, title: str, metrics: List[Dict]) -> None:
        """
        Create a slide with metric cards (KPIs).

        Args:
            title: Slide title
            metrics: List of dicts with keys: 'label', 'value', 'delta' (optional), 'color' (optional)
                     Example: {'label': 'Total Check-ins', 'value': '142', 'delta': '+12 vs last week'}
        """
        slide = self._add_blank_slide()
        self._add_title(slide, title)
        self._add_footer(slide)

        # Layout metrics in grid
        num_metrics = len(metrics)
        if num_metrics <= 2:
            cols, rows = num_metrics, 1
        elif num_metrics <= 4:
            cols, rows = 2, 2
        else:
            cols, rows = 3, 2  # Max 6 metrics

        card_width = (self.slide_width - Inches(1.5)) / cols
        card_height = (self.slide_height - Inches(2.5)) / rows
        start_top = Inches(1.5)
        start_left = Inches(0.5)

        for idx, metric in enumerate(metrics[:6]):  # Max 6 metrics
            row = idx // cols
            col = idx % cols

            left = start_left + (col * card_width) + Inches(0.1)
            top = start_top + (row * card_height) + Inches(0.1)
            width = card_width - Inches(0.2)
            height = card_height - Inches(0.2)

            # Add card background
            card = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, width, height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(240, 242, 246)  # Light gray background
            card.line.color.rgb = COLORS['gray']

            # Metric value
            value_top = top + height / 2 - Inches(0.4)
            value_text = slide.shapes.add_textbox(left, value_top, width, Inches(0.6))
            value_frame = value_text.text_frame
            value_frame.text = str(metric['value'])
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(36)
            value_para.font.bold = True
            # Handle color as string or RGBColor
            color_value = metric.get('color', 'primary_blue')
            if isinstance(color_value, str):
                color_value = COLORS.get(color_value, COLORS['primary_blue'])
            value_para.font.color.rgb = color_value
            value_para.font.name = 'Calibri'
            value_para.alignment = PP_ALIGN.CENTER

            # Label
            label_top = value_top - Inches(0.5)
            label_text = slide.shapes.add_textbox(left, label_top, width, Inches(0.4))
            label_frame = label_text.text_frame
            label_frame.text = metric['label']
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = COLORS['dark_gray']
            label_para.font.name = 'Calibri'
            label_para.alignment = PP_ALIGN.CENTER

            # Delta (optional)
            if 'delta' in metric and metric['delta']:
                delta_top = value_top + Inches(0.6)
                delta_text = slide.shapes.add_textbox(left, delta_top, width, Inches(0.3))
                delta_frame = delta_text.text_frame
                delta_frame.text = metric['delta']
                delta_para = delta_frame.paragraphs[0]
                delta_para.font.size = Pt(12)
                delta_para.font.color.rgb = COLORS['gray']
                delta_para.font.name = 'Calibri'
                delta_para.alignment = PP_ALIGN.CENTER

    def chart_slide(self, title: str, image_path: str) -> None:
        """
        Create a slide with a chart image.

        Args:
            title: Slide title
            image_path: Path to chart image file
        """
        slide = self._add_blank_slide()
        self._add_title(slide, title)
        self._add_footer(slide)

        # Add image centered
        img_left = Inches(1)
        img_top = Inches(1.8)
        img_width = self.slide_width - Inches(2)
        img_height = self.slide_height - Inches(3)

        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)

    def table_slide(self, title: str, df: pd.DataFrame, col_widths: Optional[List[Inches]] = None) -> None:
        """
        Create a slide with a data table.

        Args:
            title: Slide title
            df: Pandas DataFrame to display
            col_widths: Optional list of column widths in Inches
        """
        slide = self._add_blank_slide()
        self._add_title(slide, title)
        self._add_footer(slide)

        # Limit rows for readability
        df_display = df.head(10)
        rows, cols = len(df_display) + 1, len(df_display.columns)  # +1 for header

        # Calculate table dimensions
        table_left = Inches(0.5)
        table_top = Inches(1.8)
        table_width = self.slide_width - Inches(1)
        table_height = self.slide_height - Inches(3)

        # Add table
        table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

        # Set column widths
        if col_widths and len(col_widths) == cols:
            for idx, width in enumerate(col_widths):
                table.columns[idx].width = int(width)
        else:
            # Equal width columns
            col_width = int(table_width / cols)
            for col_idx in range(cols):
                table.columns[col_idx].width = col_width

        # Header row
        for col_idx, col_name in enumerate(df_display.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['primary_blue']
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLORS['white']
            paragraph.font.name = 'Calibri'

        # Data rows
        for row_idx, (_, row_data) in enumerate(df_display.iterrows(), start=1):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(11)
                paragraph.font.name = 'Calibri'

                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

    def bullet_points_slide(self, title: str, points: List[str], subtitle: str = "") -> None:
        """
        Create a slide with bullet points.

        Args:
            title: Slide title
            points: List of bullet point strings
            subtitle: Optional subtitle
        """
        slide = self._add_blank_slide()
        self._add_title(slide, title)
        self._add_footer(slide)

        # Add subtitle if provided
        content_top = Inches(1.8)
        if subtitle:
            self._add_text_box(
                slide, subtitle,
                Inches(0.5), Inches(1.5), self.slide_width - Inches(1), Inches(0.5),
                font_size=18, bold=True, color=COLORS['dark_gray']
            )
            content_top = Inches(2.1)

        # Add bullet points
        left = Inches(0.8)
        width = self.slide_width - Inches(1.6)
        height = self.slide_height - content_top - Inches(0.8)

        text_box = slide.shapes.add_textbox(left, content_top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for idx, point in enumerate(points):
            if idx > 0:
                text_frame.add_paragraph()

            paragraph = text_frame.paragraphs[idx]
            paragraph.text = point
            paragraph.font.size = Pt(18)
            paragraph.font.name = 'Calibri'
            paragraph.level = 0
            paragraph.space_before = Pt(12)

    def two_column_slide(self, title: str, left_content_type: str, left_content,
                         right_content_type: str, right_content) -> None:
        """
        Create a slide with two columns.

        Args:
            title: Slide title
            left_content_type: 'text', 'bullets', or 'image'
            left_content: Content for left column
            right_content_type: 'text', 'bullets', or 'image'
            right_content: Content for right column
        """
        slide = self._add_blank_slide()
        self._add_title(slide, title)
        self._add_footer(slide)

        # Column dimensions
        col_width = (self.slide_width - Inches(1.5)) / 2
        col_height = self.slide_height - Inches(2.5)
        top = Inches(1.8)
        left_col = Inches(0.5)
        right_col = left_col + col_width + Inches(0.5)

        # Left column
        self._add_column_content(slide, left_content_type, left_content, left_col, top, col_width, col_height)

        # Right column
        self._add_column_content(slide, right_content_type, right_content, right_col, top, col_width, col_height)

    def _add_column_content(self, slide, content_type: str, content, left: Inches, top: Inches,
                            width: Inches, height: Inches):
        """Helper to add content to a column."""
        if content_type == 'text':
            self._add_text_box(slide, content, left, top, width, height)
        elif content_type == 'bullets':
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            for idx, point in enumerate(content):
                if idx > 0:
                    text_frame.add_paragraph()
                paragraph = text_frame.paragraphs[idx]
                paragraph.text = point
                paragraph.font.size = Pt(16)
                paragraph.font.name = 'Calibri'
                paragraph.level = 0
        elif content_type == 'image':
            slide.shapes.add_picture(content, left, top, width=width)

    def key_takeaways_slide(self, takeaways: List[str]) -> None:
        """
        Create a "Key Takeaways" slide.

        Args:
            takeaways: List of takeaway strings
        """
        self.bullet_points_slide("🎯 Key Takeaways", takeaways)
